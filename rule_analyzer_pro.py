import re
import fitz  # PyMuPDF for PDF
import docx  # For DOCX files
import pptx  # For PPTX files
import pandas as pd
import sqlite3
import matplotlib.pyplot as plt
from concurrent.futures import ProcessPoolExecutor, as_completed
import nltk
from nltk.tokenize import sent_tokenize
from pathlib import Path

# ====================== NLTK RESOURCES ======================
nltk.download("punkt", quiet=True)
nltk.download("punkt_tab", quiet=True)

# ====================== RULE DEFINITIONS ======================
BASE_RULES = {
    "neural_network_mentions": r"\b(neural\s+network(s)?|deep\s+network(s)?|feed[- ]?forward\s+network(s)?)\b",
    "layer_terms": r"\b(layer|layers|hidden\s+layer|output\s+layer|input\s+layer|convolution(al)?\s+layer|max\s+pooling|residual\s+layer)\b",
    "neuron_terms": r"\b(neuron|neurons|relu\s+neuron|sigmoid\s+neuron|tanh\s+neuron|linear\s+neuron)\b",
    "network_operations": r"\b(forward\s+pass|backprop(agation)?|gradient\s+descent|delta\s+rule|update(d)?\s+weights?)\b",
    "architecture_terms": r"\b(autoencoder|transformer(s)?|residual\s+network|rnn(s)?|cnn(s)?|convolutional\s+network|lstm|attention)\b",
    "training_terms": r"\b(train(ed|ing)?|epoch(s)?|batch\s+size|optimizer|stochastic\s+gradient|minibatch|overfitting)\b",
    "math_terms": r"\b(matrix|matrices|vector(s)?|dot\s+product|eigenvalue(s)?|eigenvector(s)?|linear\s+algebra)\b",
    "probability_terms": r"\b(probability|random\s+variable(s)?|bayes|expectation|variance|entropy|distribution)\b",
    "evaluation_terms": r"\b(accuracy|loss|precision|recall|f1|error\s+rate|validation\s+set|test\s+set)\b",
    "compute_constraints": r"\b(scale|scalability|memory\s+constrained|computational(ly)?|expensive|efficient|optimization)\b"
}

# ====================== FILE EXTRACTION ======================
def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        doc = fitz.open(file_path)
        text = "".join(page.get_text() for page in doc)
        return text

    elif ext == ".docx":
        doc = docx.Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    elif ext == ".pptx":
        prs = pptx.Presentation(file_path)
        return "\n".join(
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        )

    else:
        raise ValueError("❌ Unsupported file type. Use PDF / DOCX / TXT / PPTX")


# ====================== TEXT PROCESSING ======================
def preprocess(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def chunk_text(text: str, size: int = 20000):
    sentences = sent_tokenize(text)
    chunks = []
    temp = ""
    length = 0

    for s in sentences:
        if length + len(s) <= size:
            temp += " " + s
            length += len(s)
        else:
            chunks.append(temp.strip())
            temp = s
            length = len(s)

    if temp:
        chunks.append(temp)

    return chunks


# ====================== RULE ENGINE ======================
def apply_rules(chunk: str, rules: dict) -> dict:
    return {name: len(re.findall(pattern, chunk, re.I)) for name, pattern in rules.items()}


def parallel_process(chunks, rules, workers: int):
    results = {}
    with ProcessPoolExecutor(max_workers=workers) as exe:
        future_map = {exe.submit(apply_rules, chunk, rules): i for i, chunk in enumerate(chunks)}
        for future in as_completed(future_map):
            idx = future_map[future]
            results[f"chunk_{idx}"] = future.result()
    return results


# ====================== SAVE + VISUALS ======================
def save_to_database(df: pd.DataFrame):
    conn = sqlite3.connect("rules.db")
    df.to_sql("results", conn, if_exists="replace", index=False)
    conn.close()


def plot_rule_totals(df: pd.DataFrame):
    numeric = df.select_dtypes(include="number")
    if numeric.empty:
        print("⚠ No numeric data to visualize.")
        return

    totals = numeric.sum()
    plt.figure()
    totals.plot(kind="bar")
    plt.title("Total Matches per Rule (Global)")
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()
    plt.show()


def plot_chunk_intensity(df: pd.DataFrame):
    numeric = df.select_dtypes(include="number")
    if numeric.empty:
        return

    df["total_matches"] = numeric.sum(axis=1)
    plt.figure()
    plt.plot(df.index, df["total_matches"])
    plt.title("Total Matches per Chunk (Document Flow)")
    plt.xlabel("Chunk Index")
    plt.ylabel("Total Matches")
    plt.tight_layout()
    plt.show()


# ====================== SUMMARY & INSIGHTS ======================
def print_keyword_summary(df: pd.DataFrame):
    numeric = df.select_dtypes(include="number")
    if numeric.empty:
        print("No keyword counts to summarize.")
        return

    totals = numeric.sum().sort_values(ascending=False)
    print("\n===== GLOBAL KEYWORD / RULE SUMMARY =====")
    for name, value in totals.items():
        print(f"{name:25s} -> {int(value)}")


def print_top_chunks(df: pd.DataFrame, top_n: int = 3, snippet_len: int = 350):
    numeric = df.select_dtypes(include="number")
    if numeric.empty:
        return

    df["total_matches"] = numeric.sum(axis=1)
    top_indices = df["total_matches"].sort_values(ascending=False).head(top_n).index

    print(f"\n===== TOP {top_n} HIGH-SIGNAL CHUNKS (by total rule matches) =====")
    for i, idx in enumerate(top_indices, start=1):
        print(f"\n--- Chunk {idx} (Rank #{i}) | Total Matches = {int(df.loc[idx, 'total_matches'])} ---")
        snippet = str(df.loc[idx, "text_chunk"])[:snippet_len].replace("\n", " ")
        print(snippet + ("..." if len(snippet) == snippet_len else ""))


# ====================== MAIN ======================
def main():
    print("\n====== DOCUMENT RULE ANALYZER PRO (LOCAL VERSION) ======\n")

    filepath = input("Enter file path to analyze: ").strip()
    if not Path(filepath).exists():
        print("❌ File not found. Try again.")
        return

    workers = int(input("Enter CPU workers (default=4): ") or 4)
    custom_terms = input("Enter custom keywords (comma separated, optional): ")

    custom_rules = {
        f"CUSTOM:{t.strip()}": rf"\b{re.escape(t.strip())}\b"
        for t in custom_terms.split(",") if t.strip()
    }

    rules = {**BASE_RULES, **custom_rules}

    print("\n[1/4] Extracting text...")
    text = preprocess(extract_text(filepath))

    print("[2/4] Splitting into chunks...")
    chunks = chunk_text(text)

    print(f"[3/4] Running parallel rule engine on {len(chunks)} chunks with {workers} workers...")
    results = parallel_process(chunks, rules, workers)

    df = pd.DataFrame(results).T
    df.insert(0, "text_chunk", chunks)

    print("[4/4] Saving results...")
    df.to_csv("output_rules.csv", index=False)
    save_to_database(df)

    print("\n✔ Analysis Complete")
    print("📁 Saved → output_rules.csv")
    print("📁 Saved → rules.db")
    print("\n===== PREVIEW (first 5 rows) =====\n")
    print(df.head(), "\n")

    # ----- Dashboard-style outputs -----
    print_keyword_summary(df)
    print_top_chunks(df, top_n=3)

    # Visuals
    print("\nOpening charts...")
    plot_rule_totals(df)
    plot_chunk_intensity(df)


if __name__ == "__main__":
    from multiprocessing import freeze_support
    freeze_support()
    main()